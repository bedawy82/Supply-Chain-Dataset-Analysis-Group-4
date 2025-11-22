# Supply_Chain_Insights_and_Forecasting.py
# Streamlit app: Supply Chain Analysis + Prophet Forecasting
# Save as:
# C:\Data Analysis\DEPI_R3\Graguation_Project\Supply_Chain_Insights_and_Forecasting.py
# Run:
# streamlit run "C:\Data Analysis\DEPI_R3\Graguation_Project\Supply_Chain_Insights_and_Forecasting.py"

import streamlit as st
import pandas as pd
import numpy as np
from pathlib import Path
from datetime import datetime
import io
import matplotlib.pyplot as plt

# Optional enhancements
try:
    import plotly.express as px
    PLOTLY_AVAILABLE = True
except Exception:
    PLOTLY_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

# Prophet import (forecasting)
try:
    from prophet import Prophet
    PROPHET_AVAILABLE = True
except Exception:
    PROPHET_AVAILABLE = False

# Page config
st.set_page_config(page_title="Supply Chain Insights & Forecasting", layout="wide", initial_sidebar_state="expanded")

# ---------------- Helper functions ----------------
@st.cache_data
def load_csv(path: Path) -> pd.DataFrame:
    """Load CSV with fallback encodings."""
    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")
    try:
        df = pd.read_csv(path, low_memory=False)
    except Exception:
        df = pd.read_csv(path, encoding="latin1", low_memory=False)
    return df

def detect_column_types(df: pd.DataFrame):
    """Detect date, numeric and categorical columns (best-effort)."""
    date_candidates = []
    for c in df.columns:
        lc = str(c).lower()
        if any(k in lc for k in ("date", "time", "order", "ship", "created", "deliv", "dispatch", "lead")):
            try:
                parsed = pd.to_datetime(df[c], errors="coerce", infer_datetime_format=True, dayfirst=True)
                if parsed.notna().sum() > 0:
                    date_candidates.append(c)
            except Exception:
                pass

    numeric_cols = [c for c in df.columns if pd.api.types.is_numeric_dtype(df[c])]
    # Try to coerce object columns to numeric if they look numeric
    for c in df.select_dtypes(include="object").columns:
        s = df[c].astype(str).str.replace(",", "").str.replace("$", "").str.replace("EGP", "", regex=False)
        coerced = pd.to_numeric(s, errors="coerce")
        if coerced.notna().sum() > 0.05 * len(df) and c not in numeric_cols:
            numeric_cols.append(c)

    categorical_cols = [c for c in df.columns if df[c].dtype == object and df[c].nunique(dropna=True) < 1000]
    return date_candidates, numeric_cols, categorical_cols

def create_excel_bytes(sheets: dict) -> io.BytesIO:
    """Create an in-memory Excel file from a dict of DataFrames."""
    bio = io.BytesIO()
    with pd.ExcelWriter(bio, engine="openpyxl") as writer:
        for name, df in sheets.items():
            safe = str(name)[:31]
            df.to_excel(writer, sheet_name=safe, index=False)
    bio.seek(0)
    return bio

def create_html_report(df: pd.DataFrame, chart_html_snippets: dict, title="Supply Chain Analysis Report"):
    """Return a simple HTML string for download (charts snippets should be HTML strings)."""
    html = "<html><head><meta charset='utf-8'><title>{}</title></head><body>".format(title)
    html += f"<h1>{title}</h1>"
    html += f"<p>Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>"
    html += "<h2>Dataset overview</h2>"
    html += f"<p>Rows: {df.shape[0]} — Columns: {df.shape[1]}</p>"
    try:
        html += "<h3>Summary (sample)</h3>"
        html += df.describe().head(10).to_html()
    except Exception:
        html += "<p>Summary unavailable.</p>"
    for k, snippet in chart_html_snippets.items():
        html += f"<h3>{k}</h3>"
        html += snippet
    html += "</body></html>"
    return html

def create_pptx_bytes(kpis: dict, image_paths: dict):
    """Create a simple PPTX in memory (requires python-pptx)."""
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx not available.")
    prs = Presentation()
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Supply Chain Insights & Forecasting"
    slide.placeholders[1].text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    # KPI slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key KPIs"
    tf = slide.shapes.placeholders[1].text_frame
    for k, v in kpis.items():
        p = tf.add_paragraph()
        p.text = f"{k}: {v}"
    # chart slides
    for title, path in image_paths.items():
        try:
            s = prs.slides.add_slide(prs.slide_layouts[5])
            s.shapes.title.text = title[:60]
            s.shapes.add_picture(str(path), Inches(1), Inches(1.2), height=Inches(4.8))
        except Exception:
            continue
    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio

def prepare_forecast_df(df, date_col, target_col):
    """Prepare DataFrame for Prophet: columns ds (date) and y (value)."""
    tmp = df[[date_col, target_col]].copy()
    tmp = tmp.dropna()
    tmp[date_col] = pd.to_datetime(tmp[date_col], errors="coerce", infer_datetime_format=True, dayfirst=True)
    tmp = tmp.dropna(subset=[date_col])
    tmp = tmp.sort_values(date_col)
    tmp = tmp.rename(columns={date_col: "ds", target_col: "y"})
    return tmp

# ---------------- Sidebar: Data selection ----------------
st.sidebar.title("Data & Pages")
st.sidebar.write("Default CSV name: `supply_chain_data.csv` (same folder as this script). You can upload a CSV.")

default_csv = Path(__file__).parent / "supply_chain_data.csv"
uploaded = st.sidebar.file_uploader("Upload CSV (optional)", type=["csv"])

if uploaded is not None:
    try:
        df = pd.read_csv(uploaded, low_memory=False)
        st.sidebar.success("Uploaded CSV loaded.")
    except Exception as e:
        st.sidebar.error(f"Failed to read uploaded CSV: {e}")
        st.stop()
else:
    if default_csv.exists():
        try:
            df = load_csv(default_csv)
            st.sidebar.success(f"Loaded default CSV: {default_csv.name}")
        except Exception as e:
            st.sidebar.error(f"Failed to load default CSV: {e}")
            st.stop()
    else:
        st.sidebar.warning("Default CSV not found. Please upload a CSV file.")
        st.stop()

# Basic cleaning
df.columns = [str(c).strip() for c in df.columns]
for c in df.select_dtypes(include="object").columns:
    df[c] = df[c].str.strip()

date_cols, numeric_cols, cat_cols = detect_column_types(df)

# Top-level page navigation (separate pages via sidebar)
page = st.sidebar.radio("Select page", ["Dashboard", "Forecasting", "Export / Reports"])

# ---------- DASHBOARD PAGE ----------
if page == "Dashboard":
    st.title("Supply Chain Insights Dashboard")
    st.markdown("Interactive exploration: filters, KPIs, charts, and pivot preview.")

    # Quick KPIs area (attempt to auto-detect revenue/cost/qty)
    kpi_candidates = {
        "Revenue": [c for c in df.columns if "revenue" in c.lower() or "sales" in c.lower() or "total" in c.lower()],
        "Cost": [c for c in df.columns if "cost" in c.lower() or "expense" in c.lower()],
        "Quantity": [c for c in df.columns if any(k in c.lower() for k in ("quantity", "qty", "units", "order_quantity"))]
    }
    kpis_display = {}
    for name, cand in kpi_candidates.items():
        total_val = None
        for col in cand:
            if pd.api.types.is_numeric_dtype(df[col]):
                try:
                    total_val = df[col].sum()
                    kpis_display[f"Total {name} ({col})"] = f"{total_val:,.2f}"
                    break
                except Exception:
                    continue
        if total_val is None:
            kpis_display[f"Total {name}"] = "N/A"

    # Show top KPIs
    cols = st.columns(len(kpis_display))
    for (k, v), c in zip(kpis_display.items(), cols):
        c.metric(k, v)

    st.markdown("---")
    # Filters in sidebar
    st.sidebar.header("Dashboard Filters")
    selected_cat = st.sidebar.selectbox("Category column (group/filter)", options=[None] + cat_cols)
    df_filtered = df.copy()
    if selected_cat:
        unique_vals = sorted(df_filtered[selected_cat].dropna().unique().tolist())
        selected_vals = st.sidebar.multiselect(f"Filter {selected_cat}", options=unique_vals[:200], default=unique_vals[:5])
        if selected_vals:
            df_filtered = df_filtered[df_filtered[selected_cat].isin(selected_vals)]

    selected_num = None
    if numeric_cols:
        selected_num = st.sidebar.selectbox("Numeric column for charts", options=[None] + numeric_cols)

    selected_date = None
    if date_cols:
        selected_date = st.sidebar.selectbox("Date column for time series", options=[None] + date_cols)
        if selected_date:
            df_filtered["_parsed_date"] = pd.to_datetime(df_filtered[selected_date], errors="coerce", infer_datetime_format=True, dayfirst=True)
            min_d = df_filtered["_parsed_date"].min()
            max_d = df_filtered["_parsed_date"].max()
            if pd.notna(min_d) and pd.notna(max_d):
                start_d, end_d = st.sidebar.date_input("Date range", value=(min_d.date(), max_d.date()))
                df_filtered = df_filtered[(df_filtered["_parsed_date"] >= pd.to_datetime(start_d)) & (df_filtered["_parsed_date"] <= pd.to_datetime(end_d))]

    # Data sample and summary
    st.subheader("Data sample")
    st.dataframe(df_filtered.head(200))

    st.subheader("Numeric summary")
    if numeric_cols:
        st.dataframe(df_filtered[numeric_cols].describe().transpose())
    else:
        st.info("No numeric columns detected.")

    st.markdown("---")
    st.subheader("Charts")

    # Time series chart
    chart_paths = {}    # to collect static images for PPTX if needed
    chart_html_snippets = {}  # for HTML report

    if selected_date and selected_num:
        st.markdown("**Time Series (daily aggregation)**")
        try:
            ts = df_filtered.dropna(subset=["_parsed_date"]).groupby(df_filtered["_parsed_date"].dt.to_period("D"))[selected_num].sum()
            ts.index = ts.index.to_timestamp()
            if PLOTLY_AVAILABLE:
                fig = px.line(ts.reset_index(), x=ts.index, y=selected_num, title=f"Daily {selected_num}")
                st.plotly_chart(fig, use_container_width=True)
                chart_html_snippets["Time Series"] = fig.to_html(full_html=False, include_plotlyjs=False)
            else:
                fig, ax = plt.subplots(figsize=(10, 4))
                ax.plot(ts.index, ts.values)
                ax.set_title(f"Daily {selected_num}")
                ax.set_ylabel(selected_num)
                ax.set_xlabel("Date")
                st.pyplot(fig)
                png = Path.cwd() / "time_series_chart.png"
                fig.savefig(png, bbox_inches='tight')
                chart_paths["Time Series"] = png
                chart_html_snippets["Time Series"] = f"<img src='{png.name}' style='max-width:800px;'/>"
        except Exception as e:
            st.warning(f"Time series error: {e}")

    # Category bar chart
    if selected_cat and selected_num:
        st.markdown(f"**Top {selected_cat} by {selected_num}**")
        try:
            agg = df_filtered.groupby(selected_cat)[selected_num].sum().sort_values(ascending=False).head(20)
            if PLOTLY_AVAILABLE:
                fig = px.bar(agg.reset_index(), x=selected_cat, y=selected_num, title=f"Top {selected_cat} by {selected_num}")
                st.plotly_chart(fig, use_container_width=True)
                chart_html_snippets["Top Category"] = fig.to_html(full_html=False, include_plotlyjs=False)
            else:
                fig, ax = plt.subplots(figsize=(10, 5))
                agg.plot(kind="bar", ax=ax)
                ax.set_ylabel(selected_num)
                ax.set_title(f"Top {selected_cat} by {selected_num}")
                st.pyplot(fig)
                png = Path.cwd() / "top_category_chart.png"
                fig.savefig(png, bbox_inches='tight')
                chart_paths["Top Category"] = png
                chart_html_snippets["Top Category"] = f"<img src='{png.name}' style='max-width:800px;'/>"
        except Exception as e:
            st.warning(f"Category chart error: {e}")

    # Distribution
    if selected_num:
        st.markdown("**Distribution**")
        try:
            if PLOTLY_AVAILABLE:
                fig = px.histogram(df_filtered, x=selected_num, nbins=50, marginal="box", title=f"Distribution of {selected_num}")
                st.plotly_chart(fig, use_container_width=True)
                chart_html_snippets["Distribution"] = fig.to_html(full_html=False, include_plotlyjs=False)
            else:
                fig, ax = plt.subplots(figsize=(8, 3))
                ax.boxplot(df_filtered[selected_num].dropna())
                ax.set_title(f"Boxplot: {selected_num}")
                st.pyplot(fig)
                png = Path.cwd() / "distribution_chart.png"
                fig.savefig(png, bbox_inches='tight')
                chart_paths["Distribution"] = png
                chart_html_snippets["Distribution"] = f"<img src='{png.name}' style='max-width:800px;'/>"
        except Exception as e:
            st.warning(f"Distribution error: {e}")

    # Pivot preview
    st.markdown("---")
    st.subheader("Pivot Preview")
    pivot_index = st.multiselect("Pivot rows (index)", options=df_filtered.columns.tolist(), default=[selected_cat] if selected_cat else [])
    pivot_value = st.selectbox("Pivot value (numeric)", options=[None] + numeric_cols)
    pivot_agg = st.selectbox("Aggregation function", options=["sum", "mean", "count", "median"], index=0)
    if pivot_index and pivot_value:
        try:
            pvt = pd.pivot_table(df_filtered, index=pivot_index, values=pivot_value, aggfunc=pivot_agg)
            st.dataframe(pvt.reset_index().head(500))
        except Exception as e:
            st.warning(f"Pivot error: {e}")

# ---------- FORECASTING PAGE ----------
elif page == "Forecasting":
    st.title("📈 Forecasting & Demand Prediction")
    st.markdown("Use Prophet to forecast a numeric target (e.g., Revenue) based on a date column.")

    if not PROPHET_AVAILABLE:
        st.error("Prophet library is not installed. Install with: `pip install prophet`. Forecasting is disabled until Prophet is available.")
        st.stop()

    # Inputs for forecasting
    st.sidebar.header("Forecast settings")
    date_col = st.sidebar.selectbox("Select date column", options=[None] + date_cols)
    target_col = st.sidebar.selectbox("Select numeric target to forecast", options=[None] + numeric_cols)
    freq = st.sidebar.selectbox("Forecast frequency", options=["D", "W", "M"], index=0)  # daily, weekly, monthly
    periods = st.sidebar.number_input("Number of future periods to forecast", min_value=1, max_value=365, value=7)
    include_history = st.sidebar.checkbox("Include historical data in plot", value=True)

    if date_col is None or target_col is None:
        st.info("Please select both a date column and a numeric target to enable forecasting.")
    else:
        # Prepare data
        df_fore = prepare_forecast_df(df, date_col, target_col)
        if df_fore.shape[0] < 10:
            st.warning("Not enough data points for a reliable forecast (need >= 10 rows after cleaning).")
        else:
            st.subheader("Training data (sample)")
            st.dataframe(df_fore.tail(200))

            # fit prophet model
            with st.spinner("Fitting Prophet model..."):
                m = Prophet()
                try:
                    m.fit(df_fore.rename(columns={"ds": "ds", "y": "y"}))
                except Exception as e:
                    st.error(f"Failed to fit Prophet: {e}")
                    st.stop()

            # predict future
            future = m.make_future_dataframe(periods=periods, freq=freq)
            forecast = m.predict(future)

            # display plot
            st.subheader("Forecast plot")
            try:
                if PLOTLY_AVAILABLE:
                    # Prophet has a built-in plotly method if using plotly; we convert to a plotly line chart
                    plot_df = forecast[['ds', 'yhat', 'yhat_lower', 'yhat_upper']].set_index('ds')
                    if include_history:
                        # overlay actuals
                        actuals = df_fore.set_index('ds')['y']
                        plot_df = plot_df.join(actuals.rename('actual'), how='left')
                        fig = px.line(plot_df.reset_index(), x='ds', y=['actual', 'yhat'], title=f"Forecast for {target_col}")
                        # add bands
                        fig.add_traces([
                            px.line(plot_df.reset_index(), x='ds', y='yhat_lower').data[0],
                            px.line(plot_df.reset_index(), x='ds', y='yhat_upper').data[0]
                        ])
                        st.plotly_chart(fig, use_container_width=True)
                        # store html snippet
                        chart_html = fig.to_html(full_html=False, include_plotlyjs=False)
                        chart_html_snippets = {"Forecast": chart_html}
                    else:
                        fig = px.line(plot_df.reset_index(), x='ds', y='yhat', title=f"Forecast for {target_col}")
                        st.plotly_chart(fig, use_container_width=True)
                        chart_html_snippets = {"Forecast": fig.to_html(full_html=False, include_plotlyjs=False)}
                else:
                    fig = m.plot(forecast)
                    st.pyplot(fig)
                    # Save static image for PPTX
                    png = Path.cwd() / "prophet_forecast.png"
                    fig.savefig(png, bbox_inches='tight')
                    chart_paths["Forecast"] = png
                    chart_html_snippets = {"Forecast": f"<img src='{png.name}' style='max-width:800px;'/>"}
            except Exception as e:
                st.warning(f"Plotting forecast failed: {e}")
                chart_html_snippets = {}

            # show forecast tail table
            st.subheader("Forecast table (tail)")
            show_cols = ["ds", "yhat", "yhat_lower", "yhat_upper"]
            st.dataframe(forecast[show_cols].tail(periods))

            # download forecast as excel
            if st.button("Download forecast as Excel"):
                try:
                    sheets = {"Forecast": forecast[show_cols].reset_index(drop=True)}
                    sheets["Train"] = df_fore.rename(columns={"ds": date_col, "y": target_col})
                    excel_io = create_excel_bytes(sheets)
                    fn = f"forecast_{target_col}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
                    st.download_button("Click to download Excel", data=excel_io, file_name=fn, mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
                except Exception as e:
                    st.error(f"Failed to prepare Excel: {e}")

# ---------- EXPORT / REPORTS PAGE ----------
elif page == "Export / Reports":
    st.title("Export & Reports")
    st.markdown("Generate Excel summaries, HTML report, or PPTX slides.")

    st.sidebar.header("Export options")
    if st.button("Generate Excel summary (Raw + Summary + Top categories)"):
        try:
            sheets = {}
            sheets["Raw_Data"] = df.copy()
            try:
                sheets["Summary"] = df.describe(include="all").transpose().reset_index()
            except Exception:
                sheets["Summary"] = pd.DataFrame()
            # top values for first 6 categorical columns
            ct = 0
            for c in cat_cols:
                if ct >= 6:
                    break
                tv = df[c].value_counts(dropna=False).rename_axis(c).reset_index(name="count")
                sheets[f"Top_{c[:20]}"] = tv
                ct += 1
            excel_io = create_excel_bytes(sheets)
            fn = f"supply_chain_summary_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
            st.download_button("Download Excel", data=excel_io, file_name=fn, mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        except Exception as e:
            st.error(f"Failed to generate Excel: {e}")

    if st.button("Generate HTML report"):
        try:
            # prepare small chart snippets dictionary if available from previous pages (best-effort)
            chart_snippets = {}
            html = create_html_report(df, chart_snippets, title="Supply Chain Analysis Report")
            b = html.encode("utf-8")
            hfn = f"supply_chain_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
            st.download_button("Download HTML report", data=b, file_name=hfn, mime="text/html")
        except Exception as e:
            st.error(f"Failed to create HTML report: {e}")

    if st.button("Generate PPTX summary"):
        if not PPTX_AVAILABLE:
            st.error("python-pptx not installed. Install via: pip install python-pptx")
        else:
            try:
                # prepare minimal KPIs
                kpis_for_ppt = {k: v for k, v in kpis_display.items()} if 'kpis_display' in globals() else {}
                # use PNGs created earlier if any in chart_paths
                pptx_bio = create_pptx_bytes(kpis_for_ppt, {k:str(v) for k,v in chart_paths.items()} if 'chart_paths' in globals() else {})
                pf = f"supply_chain_slides_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
                st.download_button("Download PPTX", data=pptx_bio, file_name=pf, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
            except Exception as e:
                st.error(f"Failed to create PPTX: {e}")

st.markdown("---")
st.caption("Supply Chain Insights & Forecasting — Streamlit app. If you want custom KPIs or automated scheduling/exporting, tell me the KPI column names and desired schedule.")
